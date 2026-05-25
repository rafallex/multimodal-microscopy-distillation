"""Build A3_cancer_challenge.pptx — 11-slide deck for the Kaggle challenge presentation.

Run from the project root or this directory; output lands next to this script.

Palette: Midnight Executive variant
- NAVY  #1A2540  (title-slide bg, dark text)
- SLATE #2E5266  (section bars)
- TERRA #E07856  (accent / wins)
- CREAM #F5F0E8  (body bg)
- MUTED #8A9197  (secondary text)

Typography: Georgia (titles, 36-44pt) + Calibri (body, 14-16pt).
Motif: 0.08" terracotta accent bar at left edge of every content slide.
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# --- Palette ---
NAVY = RGBColor(0x1A, 0x25, 0x40)
SLATE = RGBColor(0x2E, 0x52, 0x66)
TERRA = RGBColor(0xE0, 0x78, 0x56)
CREAM = RGBColor(0xF5, 0xF0, 0xE8)
MUTED = RGBColor(0x8A, 0x91, 0x97)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x1A, 0x25, 0x40)

HERE = Path(__file__).parent.resolve()
PROJECT_ROOT = HERE.parent
OUT_PATH = HERE / "A3_cancer_challenge.pptx"
LEARNING_CURVES_PNG = PROJECT_ROOT / "results" / "v19" / "learning_curves.png"


# --- Helpers ---
def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, fill_color, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
    shape.shadow.inherit = False
    return shape


def add_text(slide, x, y, w, h, text, *,
             font="Calibri", size=14, bold=False, italic=False,
             color=DARK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             line_spacing=None):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    # Single-paragraph: split text on \n
    lines = text.split("\n") if isinstance(text, str) else [text]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_spacing is not None:
            p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
    return tb


def add_bullets(slide, x, y, w, h, items, *,
                font="Calibri", size=14, color=DARK,
                bullet_color=None, line_spacing=1.15):
    """items: list of (text, [optional sub_items_list]) or just list of strings.

    Renders manual "—" bullets (cleaner than PPT auto-bullets via XML hacks)."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    bc = bullet_color if bullet_color else TERRA
    first = True
    for item in items:
        if isinstance(item, tuple):
            if len(item) == 1:
                text, subs = item[0], []
            else:
                text, subs = item[0], item[1]
        else:
            text, subs = item, []
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        run_b = p.add_run()
        run_b.text = "— "
        run_b.font.name = font
        run_b.font.size = Pt(size)
        run_b.font.bold = True
        run_b.font.color.rgb = bc
        run_t = p.add_run()
        run_t.text = text
        run_t.font.name = font
        run_t.font.size = Pt(size)
        run_t.font.color.rgb = color
        for sub in subs:
            p2 = tf.add_paragraph()
            p2.alignment = PP_ALIGN.LEFT
            p2.line_spacing = line_spacing
            p2.level = 1
            spacer = p2.add_run()
            spacer.text = "        • "
            spacer.font.name = font
            spacer.font.size = Pt(size - 2)
            spacer.font.color.rgb = MUTED
            run_s = p2.add_run()
            run_s.text = sub
            run_s.font.name = font
            run_s.font.size = Pt(size - 2)
            run_s.font.color.rgb = MUTED
    return tb


def add_accent_bar(slide):
    """0.08" terracotta vertical bar at left edge — visual motif on every content slide."""
    add_rect(slide, Inches(0), Inches(0), Inches(0.08), Inches(5.625), TERRA)


def add_footer(slide, page_num, total=12):
    add_text(slide, Inches(0.4), Inches(5.30), Inches(6), Inches(0.25),
             "Multimodal Cancer Cell Classification  ·  Uppsala 1MD042 A3",
             font="Calibri", size=9, color=MUTED, align=PP_ALIGN.LEFT)
    add_text(slide, Inches(8.5), Inches(5.30), Inches(1.2), Inches(0.25),
             f"{page_num} / {total}",
             font="Calibri", size=9, color=MUTED, align=PP_ALIGN.RIGHT)


def add_title(slide, title, eyebrow=None, y=0.35):
    if eyebrow:
        add_text(slide, Inches(0.4), Inches(y), Inches(9), Inches(0.3),
                 eyebrow.upper(),
                 font="Calibri", size=11, bold=True, color=TERRA,
                 align=PP_ALIGN.LEFT)
        y += 0.32
    add_text(slide, Inches(0.4), Inches(y), Inches(9), Inches(0.7),
             title, font="Georgia", size=28, bold=True, color=DARK,
             align=PP_ALIGN.LEFT)


def set_notes(slide, notes_text):
    notes_tf = slide.notes_slide.notes_text_frame
    notes_tf.text = notes_text


# --- Build presentation ---
pres = Presentation()
pres.slide_width = Inches(10)
pres.slide_height = Inches(5.625)
BLANK = pres.slide_layouts[6]


# =========================
# Slide 1: TITLE
# =========================
s = pres.slides.add_slide(BLANK)
set_bg(s, NAVY)

# Terracotta bar on left
add_rect(s, Inches(0), Inches(0), Inches(0.25), Inches(5.625), TERRA)

# Eyebrow
add_text(s, Inches(0.8), Inches(0.9), Inches(8), Inches(0.3),
         "ASSIGNMENT 3  ·  KAGGLE CHALLENGE",
         font="Calibri", size=12, bold=True, color=TERRA,
         align=PP_ALIGN.LEFT)

# Title
add_text(s, Inches(0.8), Inches(1.35), Inches(8.5), Inches(1.5),
         "Multimodal Cancer Cell\nClassification",
         font="Georgia", size=40, bold=True, color=WHITE,
         align=PP_ALIGN.LEFT, line_spacing=1.05)

# Subtitle
add_text(s, Inches(0.8), Inches(3.05), Inches(8.5), Inches(0.5),
         "From SSL collapse to #1 on the leaderboard — 30 versions, 5 acts, +0.081 LB",
         font="Georgia", size=18, italic=True, color=CREAM,
         align=PP_ALIGN.LEFT)

# Divider line
add_rect(s, Inches(0.8), Inches(3.85), Inches(2.5), Inches(0.025), TERRA)

# Author / course block
add_text(s, Inches(0.8), Inches(4.05), Inches(8.5), Inches(0.3),
         "Rafael Tavares Proença",
         font="Calibri", size=14, bold=True, color=WHITE,
         align=PP_ALIGN.LEFT)
add_text(s, Inches(0.8), Inches(4.40), Inches(8.5), Inches(0.3),
         "Advanced Deep Learning for Image Processing (1MD042)",
         font="Calibri", size=12, color=CREAM, align=PP_ALIGN.LEFT)
add_text(s, Inches(0.8), Inches(4.70), Inches(8.5), Inches(0.3),
         "Uppsala University  ·  June 4, 2026",
         font="Calibri", size=12, color=MUTED, align=PP_ALIGN.LEFT)

set_notes(s, (
    "Hello — I'm Rafael, presenting our run at the Multimodal Cancer Cell "
    "Classification Challenge.\n\n"
    "Quick framing: this talk isn't 'here's the model that worked'. It's the "
    "story of 36 versions and the lesson each one taught us about the gap "
    "between cross-validation AUC and public leaderboard AUC. The arc goes from "
    "an ambitious SSL pipeline that scored highest CV ever but collapsed on the LB, "
    "through a careful but unsuccessful 9-change rescue attempt, to a much "
    "simpler EfficientNet recipe at LB 0.7455 — which then became the foundation for a "
    "semi-supervised pipeline (Lee 2013 pseudo-labels, then Hinton 2015 distillation) that "
    "reached LB 0.8236 and currently holds #1 on the public leaderboard.\n\n"
    "Time budget: ~50 seconds per slide, 11 slides, leaves ~30s of buffer for Q&A."
))


# =========================
# Slide 2: THE CHALLENGE
# =========================
s = pres.slides.add_slide(BLANK)
set_bg(s, CREAM)
add_accent_bar(s)
add_title(s, "Binary classification on paired microscopy",
          eyebrow="The challenge")

# Left column: bullets
add_bullets(s, Inches(0.4), Inches(1.5), Inches(5.0), Inches(3.5), [
    ("Task: per-cell label — malignant vs benign (oral cancer)",),
    ("Inputs: paired BF + FL microscopy, 128×128 grayscale",),
    ("Train: 12 patients with leave-one-patient-out CV",),
    ("Test: unknown patients → OOD generalization is the dominant failure mode",),
    ("Current LB position: we hold #1 at 0.8264 (+0.013 above next-best)",),
], size=14)

# Right column: stat cards
def stat_card(slide, x, y, w, h, big, big_color, small, big_size=32, small_y=None):
    add_rect(slide, x, y, w, h, WHITE)
    add_rect(slide, x, y, Inches(0.06), h, TERRA)
    add_text(slide, x + Inches(0.25), y + Inches(0.12), w - Inches(0.35), Inches(0.7),
             big, font="Georgia", size=big_size, bold=True, color=big_color,
             align=PP_ALIGN.LEFT)
    if small:  # don't render an empty placeholder textbox
        sy = small_y if small_y is not None else 0.85
        add_text(slide, x + Inches(0.25), y + Inches(sy), w - Inches(0.35), Inches(0.30),
                 small, font="Calibri", size=10, color=MUTED, align=PP_ALIGN.LEFT)

stat_card(s, Inches(5.8), Inches(1.5), Inches(3.7), Inches(1.30),
          "12 patients", DARK, "training set size")
stat_card(s, Inches(5.8), Inches(2.95), Inches(3.7), Inches(1.30),
          "2 modalities", DARK, "brightfield + fluorescence per cell")
# Two small cards on bottom row: short labels rendered as `small` inside the card
stat_card(s, Inches(5.8), Inches(4.40), Inches(1.78), Inches(0.78),
          "128×", DARK, "input pixel side", big_size=26, small_y=0.48)
stat_card(s, Inches(7.72), Inches(4.40), Inches(1.78), Inches(0.78),
          "OOD", TERRA, "unseen test patients", big_size=26, small_y=0.48)

add_footer(s, 2)
set_notes(s, (
    "The setup. Per-cell binary classification — is this oral cell malignant or not — "
    "from a pair of microscopy images: brightfield and fluorescence, both 128 by 128 "
    "grayscale.\n\n"
    "Two things make this hard. First, only 12 training patients — small dataset. Second, "
    "test patients are unknown. That second point ends up dominating everything. Our cross-"
    "validation has to be leave-one-patient-out, and even that turns out not to be enough.\n\n"
    "Anchor for the audience: by the end of this presentation we currently hold the "
    "public LB at 0.8264 — but I'll trace how we got there through 5 acts. The earlier "
    "leader was at 0.78; we passed it with v46's distillation step and extended further with v47."
))


# =========================
# Slide 3: JOURNEY OVERVIEW
# =========================
s = pres.slides.add_slide(BLANK)
set_bg(s, CREAM)
add_accent_bar(s)
add_title(s, "36 versions, 4 acts", eyebrow="The journey")

# Timeline track — centered on slide, symmetric phase markers
TIMELINE_Y = Inches(2.7)
TIMELINE_X = Inches(1.5)
TIMELINE_W = Inches(7.0)
add_rect(s, TIMELINE_X, TIMELINE_Y, TIMELINE_W, Inches(0.06), MUTED)

# Four phase markers along the line (px symmetric around slide centre 5.0)
phases = [
    (1.70, "Act 1+2", "v10–v16", "Baselines & SSL", TERRA),
    (3.90, "Act 3", "v17–v19", "EffNet pivot", SLATE),
    (6.10, "Act 4", "v20–v41", "Regularizer search", DARK),
    (8.30, "Act 5", "v42–v47", "Semi-supervised", RGBColor(0x2A, 0x8A, 0x4A)),
]
for px, act, vers, label, color in phases:
    # Dot
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL,
                             Inches(px - 0.12), TIMELINE_Y - Inches(0.10),
                             Inches(0.30), Inches(0.30))
    dot.fill.solid()
    dot.fill.fore_color.rgb = color
    dot.line.fill.background()
    # Labels centered on px; width 1.6 means box spans px-0.8 .. px+0.8
    add_text(s, Inches(px - 0.8), Inches(1.55), Inches(1.6), Inches(0.3),
             act.upper(), font="Calibri", size=10, bold=True,
             color=TERRA, align=PP_ALIGN.CENTER)
    add_text(s, Inches(px - 0.8), Inches(1.85), Inches(1.6), Inches(0.3),
             vers, font="Georgia", size=13, bold=True, color=DARK,
             align=PP_ALIGN.CENTER)
    add_text(s, Inches(px - 0.95), Inches(2.15), Inches(1.9), Inches(0.4),
             label, font="Calibri", size=11, italic=True, color=MUTED,
             align=PP_ALIGN.CENTER)

# Bottom row: result per act (same px positions, narrower boxes)
results = [
    (1.70, "0.572 → 0.503", "SSL collapse"),
    (3.90, "0.7455", "supervised floor"),
    (6.10, "0.7563", "v41 stacked L4 regs"),
    (8.30, "0.8264", "v47 #1 on LB"),
]
for px, lb, note in results:
    add_text(s, Inches(px - 0.95), Inches(3.20), Inches(1.9), Inches(0.4),
             lb, font="Georgia", size=15, bold=True, color=DARK,
             align=PP_ALIGN.CENTER)
    add_text(s, Inches(px - 0.95), Inches(3.60), Inches(1.9), Inches(0.4),
             note, font="Calibri", size=10, italic=True, color=MUTED,
             align=PP_ALIGN.CENTER)

# Bottom banner
add_rect(s, Inches(0.4), Inches(4.40), Inches(9.2), Inches(0.65), NAVY)
add_text(s, Inches(0.6), Inches(4.50), Inches(9.0), Inches(0.45),
         "Each act answered a question. Each result changed the next question.",
         font="Georgia", size=14, italic=True, color=CREAM,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)

add_footer(s, 3)
set_notes(s, (
    "Four acts to this story. I'm using version numbers as anchors — we tracked "
    "every Kaggle submission in a LB_HISTORY.md file in the repo, so 'v19' or 'v34' "
    "means a specific notebook with a specific recipe.\n\n"
    "Act 1, v10 through v14, was iteration on a course-starter ResNet-18 — AdamW, "
    "OneCycleLR, paired augmentation. No LB submissions yet.\n\n"
    "Act 2, v15 and v16, was where we went big — CoMIR-style contrastive SSL. v15 "
    "scored 0.866 in CV, the highest we'd ever seen, then crashed to 0.572 on the "
    "leaderboard. v16 was a careful 9-change response that improved patient-level OOF "
    "to 0.943 but the LB stayed near random at 0.503.\n\n"
    "Act 3, v17–v19, was the pivot — drop the SSL complexity, go back to a simple "
    "EffNet-B0 + heavy augmentation. v19 landed our best single-model LB at 0.7455.\n\n"
    "Act 4, v20–v41, was about understanding what each ingredient actually "
    "contributed. Lots of failures with lessons. v34 — ResNet-50 cleanly — hit 0.7155. "
    "Act 4 climaxed in v41, which stacked four L4 regularizers cleanly for +0.011 over v19 → 0.7563.\n\n"
    "Act 5 — the new chapter, v42 through v47 — was the semi-supervised pivot. "
    "v44 added Lee 2013 hard pseudo-labels (+0.025 LB), v46 swapped to Hinton 2015 "
    "soft-target distillation on all 59k test cells (+0.039 LB), and v47 iterated to "
    "round 2 with v46 as the new teacher (+0.003 LB, Xie 2020). v47 at 0.8264 holds #1."
))


# =========================
# Slide 4: SSL AMBITION + OOD SHOCK
# =========================
s = pres.slides.add_slide(BLANK)
set_bg(s, CREAM)
add_accent_bar(s)
add_title(s, "v15: CoMIR SSL — the CV was beautiful",
          eyebrow="Act 2  ·  the ambition")

# Left: method bullets
add_text(s, Inches(0.4), Inches(1.45), Inches(5.0), Inches(0.35),
         "What we built", font="Calibri", size=12, bold=True,
         color=TERRA, align=PP_ALIGN.LEFT)
add_bullets(s, Inches(0.4), Inches(1.80), Inches(5.0), Inches(2.5), [
    ("Two ResNet-18 branches (BF + FL), each with projection head",),
    ("CoMIR contrastive SSL pretrain (NT-Xent, τ=0.1)",),
    ("Supervised stage with discriminative LR (1:10 backbone:head)",),
    ("3-fold stratified-group CV — our undoing",),
], size=13)

# Right: the shock — two big stat cards stacked
add_text(s, Inches(5.7), Inches(1.45), Inches(4.0), Inches(0.35),
         "Then the leaderboard arrived", font="Calibri", size=12, bold=True,
         color=TERRA, align=PP_ALIGN.LEFT)

# CV card
add_rect(s, Inches(5.7), Inches(1.85), Inches(3.8), Inches(1.30), WHITE)
add_rect(s, Inches(5.7), Inches(1.85), Inches(0.08), Inches(1.30), SLATE)
add_text(s, Inches(5.9), Inches(1.95), Inches(3.6), Inches(0.4),
         "CV AUC", font="Calibri", size=11, color=MUTED, align=PP_ALIGN.LEFT)
add_text(s, Inches(5.9), Inches(2.25), Inches(3.6), Inches(0.8),
         "0.866", font="Georgia", size=44, bold=True, color=DARK,
         align=PP_ALIGN.LEFT)

# LB card
add_rect(s, Inches(5.7), Inches(3.25), Inches(3.8), Inches(1.30), WHITE)
add_rect(s, Inches(5.7), Inches(3.25), Inches(0.08), Inches(1.30), TERRA)
add_text(s, Inches(5.9), Inches(3.35), Inches(3.6), Inches(0.4),
         "Public LB AUC", font="Calibri", size=11, color=MUTED,
         align=PP_ALIGN.LEFT)
add_text(s, Inches(5.9), Inches(3.65), Inches(3.6), Inches(0.8),
         "0.572", font="Georgia", size=44, bold=True, color=TERRA,
         align=PP_ALIGN.LEFT)

# Diagnosis strip
add_rect(s, Inches(0.4), Inches(4.55), Inches(9.2), Inches(0.55), NAVY)
add_text(s, Inches(0.6), Inches(4.62), Inches(9.0), Inches(0.4),
         "Diagnosis: 3-fold CV had val patients similar to train in 2/3 folds. "
         "The hard fold peaked at epoch 0 — essentially untrained.",
         font="Georgia", size=12, italic=True, color=CREAM,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)

add_footer(s, 4)
set_notes(s, (
    "v15 was our big swing. We built a CoMIR-style contrastive SSL pipeline — two "
    "ResNet-18 backbones, one for brightfield and one for fluorescence, each with a "
    "projection head trained with NT-Xent contrastive loss. The idea was to learn "
    "cross-modal correspondence first, then fine-tune for the classification task.\n\n"
    "Cross-validation came in at 0.866 — the best we'd ever seen. We were excited. "
    "We submitted.\n\n"
    "Leaderboard came in at 0.572. Twenty-nine AUC points of generalization gap.\n\n"
    "The diagnosis took us a couple of days. The 3-fold CV split was hiding it: two of "
    "three folds happened to have validation patients that were similar to training "
    "patients — val AUC > 0.88, looked great. The third fold had OOD-like patients, "
    "with val AUC peaking at epoch 0 — meaning the model never actually learned "
    "anything for that fold; it was essentially the initialization plus one gradient "
    "step. The CV mean of 0.866 buried the failure."
))


# =========================
# Slide 5: v16 RESPONSE
# =========================
s = pres.slides.add_slide(BLANK)
set_bg(s, CREAM)
add_accent_bar(s)
add_title(s, "v16: nine careful fixes",
          eyebrow="Act 2  ·  the response")

# Left: 9 changes in a 3x3 grid of small cards
changes = [
    ("1", "No mixup",       "fought CoMIR alignment"),
    ("2", "LOPO 12-fold CV","honest OOD estimator"),
    ("3", "Heavy stain aug","invariance during SSL"),
    ("4", "Logit-space TTA","preserves range"),
    ("5", "Multi-snapshot", "robust to calib drift"),
    ("6", "Label smoothing","ε = 0.05"),
    ("7", "Snap at {3, 5}", "no val_auc selection"),
    ("8", "EMA decay 0.99", "tuned for 426 steps"),
    ("9", "Aux NT-Xent",    "anchor to SSL features"),
]
grid_x0, grid_y0 = 0.4, 1.45
cw, ch = 1.85, 0.95
gx, gy = 0.05, 0.10
for idx, (num, name, why) in enumerate(changes):
    row, col = idx // 3, idx % 3
    x = Inches(grid_x0 + col * (cw + gx))
    y = Inches(grid_y0 + row * (ch + gy))
    add_rect(s, x, y, Inches(cw), Inches(ch), WHITE)
    add_rect(s, x, y, Inches(0.05), Inches(ch), SLATE)
    add_text(s, x + Inches(0.15), y + Inches(0.05), Inches(0.3), Inches(0.3),
             num, font="Georgia", size=14, bold=True, color=TERRA,
             align=PP_ALIGN.LEFT)
    add_text(s, x + Inches(0.50), y + Inches(0.05), Inches(cw - 0.55), Inches(0.35),
             name, font="Calibri", size=11, bold=True, color=DARK,
             align=PP_ALIGN.LEFT)
    add_text(s, x + Inches(0.15), y + Inches(0.42), Inches(cw - 0.2), Inches(0.5),
             why, font="Calibri", size=9, italic=True, color=MUTED,
             align=PP_ALIGN.LEFT)

# Right column: outcome
add_text(s, Inches(6.7), Inches(1.45), Inches(3.0), Inches(0.35),
         "Outcome", font="Calibri", size=12, bold=True, color=TERRA,
         align=PP_ALIGN.LEFT)

add_rect(s, Inches(6.7), Inches(1.80), Inches(3.0), Inches(1.20), WHITE)
add_rect(s, Inches(6.7), Inches(1.80), Inches(0.06), Inches(1.20), SLATE)
add_text(s, Inches(6.85), Inches(1.87), Inches(2.8), Inches(0.3),
         "Patient-level OOF", font="Calibri", size=10, color=MUTED)
add_text(s, Inches(6.85), Inches(2.13), Inches(2.8), Inches(0.7),
         "0.943", font="Georgia", size=32, bold=True, color=DARK)
add_text(s, Inches(6.85), Inches(2.65), Inches(2.8), Inches(0.3),
         "up from v15's 0.914",
         font="Calibri", size=9, italic=True, color=MUTED)

add_rect(s, Inches(6.7), Inches(3.10), Inches(3.0), Inches(1.20), WHITE)
add_rect(s, Inches(6.7), Inches(3.10), Inches(0.06), Inches(1.20), TERRA)
add_text(s, Inches(6.85), Inches(3.17), Inches(2.8), Inches(0.3),
         "Public LB", font="Calibri", size=10, color=MUTED)
add_text(s, Inches(6.85), Inches(3.43), Inches(2.8), Inches(0.7),
         "0.503", font="Georgia", size=32, bold=True, color=TERRA)
add_text(s, Inches(6.85), Inches(3.95), Inches(2.8), Inches(0.3),
         "near-random — LOPO can't see this OOD",
         font="Calibri", size=9, italic=True, color=MUTED)

# Bottom strip
add_rect(s, Inches(0.4), Inches(4.55), Inches(9.2), Inches(0.55), NAVY)
add_text(s, Inches(0.6), Inches(4.62), Inches(9.0), Inches(0.4),
         "The remaining gap is site/protocol-level OOD: every train patient shares "
         "something no test patient has.",
         font="Georgia", size=12, italic=True, color=CREAM,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)

add_footer(s, 5)
set_notes(s, (
    "We responded with v16. Nine changes, each mapped to a specific cause of v15's "
    "collapse. The big four were: no mixup (it fought the SSL alignment), LOPO CV for "
    "an honest OOD estimator, an auxiliary NT-Xent loss to anchor the supervised "
    "stage to the SSL features, and a multi-snapshot ensemble at epochs 3 and 5 instead "
    "of best-val-AUC selection. The other five were defensible smaller fixes — label "
    "smoothing, EMA tuned for our short training, heavy stain augmentation during SSL.\n\n"
    "Patient-level OOF went up to 0.943, our best ever. We thought we had it.\n\n"
    "Leaderboard came in at 0.503. Basically random.\n\n"
    "The lesson: there's a level of OOD that LOPO cannot detect. If every training "
    "patient was scanned on the same microscope or stained with the same protocol, "
    "and the test patients weren't, then no training-only validation can warn you. "
    "We'd hit a ceiling that the architecture couldn't break through."
))


# =========================
# Slide 6: THE PIVOT TO v19
# =========================
s = pres.slides.add_slide(BLANK)
set_bg(s, CREAM)
add_accent_bar(s)
add_title(s, "v19: simpler model, test-time tricks",
          eyebrow="Act 3  ·  the pivot")

# Left: 5 ingredients
add_text(s, Inches(0.4), Inches(1.40), Inches(5.0), Inches(0.35),
         "The recipe", font="Calibri", size=12, bold=True, color=TERRA)
add_bullets(s, Inches(0.4), Inches(1.75), Inches(5.5), Inches(3.2), [
    ("EfficientNet-B0 backbone (single, ImageNet-pretrained)",),
    ("Multiple-Instance Learning aggregation across cells per patient",),
    ("Heavy color/stain augmentation — ColorJitter + RandomGamma",),
    ("Adaptive BatchNorm to the test distribution",),
    ("8-way D4 Test-Time Augmentation (4 rotations × 2 reflections)",),
    ("Test-time stain normalization",),
], size=13)

# Right: massive LB callout
add_rect(s, Inches(6.4), Inches(1.40), Inches(3.2), Inches(2.85), NAVY)
add_text(s, Inches(6.55), Inches(1.55), Inches(2.9), Inches(0.35),
         "PUBLIC LB", font="Calibri", size=11, bold=True, color=TERRA,
         align=PP_ALIGN.LEFT)
add_text(s, Inches(6.55), Inches(1.95), Inches(2.9), Inches(1.2),
         "0.7455", font="Georgia", size=58, bold=True, color=WHITE,
         align=PP_ALIGN.LEFT)
add_text(s, Inches(6.55), Inches(3.15), Inches(2.9), Inches(0.4),
         "best single model",
         font="Calibri", size=12, italic=True, color=CREAM, align=PP_ALIGN.LEFT)
add_text(s, Inches(6.55), Inches(3.55), Inches(2.9), Inches(0.4),
         "+0.24 vs v16",
         font="Calibri", size=11, color=TERRA, align=PP_ALIGN.LEFT)
add_text(s, Inches(6.55), Inches(3.85), Inches(2.9), Inches(0.3),
         "vs leader 0.7832",
         font="Calibri", size=10, color=MUTED, align=PP_ALIGN.LEFT)

# Bottom insight strip
add_rect(s, Inches(0.4), Inches(4.55), Inches(9.2), Inches(0.55), NAVY)
add_text(s, Inches(0.6), Inches(4.62), Inches(9.0), Inches(0.4),
         "The win wasn't a new architecture. It was treating the test distribution "
         "as the thing to model.",
         font="Georgia", size=12, italic=True, color=CREAM,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)

add_footer(s, 6)
set_notes(s, (
    "The pivot. We dropped the SSL complexity entirely and went back to something "
    "much simpler: an off-the-shelf EfficientNet-B0, ImageNet-pretrained.\n\n"
    "What changed was where we spent the effort. Instead of trying to learn a better "
    "internal representation, we focused on test-time techniques that close the "
    "training-to-test distribution gap. Five ingredients: heavy color and stain "
    "augmentation, multiple-instance learning across cells per patient, adaptive "
    "BatchNorm — letting BN re-estimate its statistics on the test set — 8-way "
    "D4 test-time augmentation, and test-time stain normalization.\n\n"
    "Public LB came in at 0.7455. That's a 0.24 AUC jump over v16. Best single "
    "model we'd produced, and within 0.04 of the leaderboard leader.\n\n"
    "The key insight, written on the bottom of the slide: the win came from treating "
    "the test distribution as the thing to model. The SSL pipeline was trying to "
    "learn invariance from scratch; v19 imported it."
))


# =========================
# Slide 7: WHAT DIDN'T WORK
# =========================
s = pres.slides.add_slide(BLANK)
set_bg(s, CREAM)
add_accent_bar(s)
add_title(s, "What didn't work — and what it taught us",
          eyebrow="Act 4  ·  the failures")

# Table-like layout: 6 failure cards in 2 columns x 3 rows
failures = [
    ("v20", "0.5974", "New aug recipe overshot — destroyed the signal."),
    ("v21", "0.7018", "ResNet-50 with 224 upscale. Confounded: upscale was the problem, not the backbone."),
    ("v22", "0.7422", "Ensemble of v19 + v21. Weaker member dragged the mean below v19 alone."),
    ("v23", "0.7154", "Same v19 recipe, different seed. 0.03 LB shift from seed alone."),
    ("v27", "—",      "Single 2-patient val holdout on N=12 → best_epoch=0 collapse."),
    ("v30", "0.6448", "Discriminative LR at 0.1 ratio. ImageNet → microscopy needs MORE backbone LR, not less."),
]
cw, ch = 4.45, 1.00
gx, gy = 0.10, 0.10
for idx, (ver, lb, lesson) in enumerate(failures):
    row, col = idx // 2, idx % 2
    x = Inches(0.4 + col * (cw + gx))
    y = Inches(1.40 + row * (ch + gy))
    add_rect(s, x, y, Inches(cw), Inches(ch), WHITE)
    add_rect(s, x, y, Inches(0.06), Inches(ch), TERRA)
    # Version
    add_text(s, x + Inches(0.20), y + Inches(0.10), Inches(0.8), Inches(0.4),
             ver, font="Georgia", size=18, bold=True, color=DARK,
             align=PP_ALIGN.LEFT)
    # LB
    add_text(s, x + Inches(0.20), y + Inches(0.50), Inches(0.9), Inches(0.35),
             "LB " + lb, font="Calibri", size=10, color=MUTED, align=PP_ALIGN.LEFT)
    # Lesson
    add_text(s, x + Inches(1.15), y + Inches(0.15), Inches(cw - 1.30), Inches(0.75),
             lesson, font="Calibri", size=11, color=DARK, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.MIDDLE)

add_footer(s, 7)
set_notes(s, (
    "Act 4 was about understanding what each ingredient in v19 actually contributed, "
    "by trying to change one thing at a time. Many of these failed. Each failure "
    "taught us something.\n\n"
    "v20: a more aggressive augmentation recipe destroyed the signal — 0.5974. "
    "There's a ceiling on how much you can perturb a 128x128 cell image before it "
    "stops being a cell.\n\n"
    "v21: tried ResNet-50, but with 224 upscale. Two changes at once. LB 0.7018. "
    "Couldn't tell whether ResNet-50 was worse or the upscaling was the problem. "
    "v34 later resolves this.\n\n"
    "v22: ensembled v19 with v21. Got 0.7422 — worse than v19 alone. Lesson: don't "
    "ensemble when one member is more than 0.02 below the other.\n\n"
    "v23: same v19 recipe, only the random seed changed. 0.7154 instead of 0.7455. "
    "0.03 LB swing from seed alone. This is the most important slide-7 lesson: "
    "any conclusion based on a single-seed delta smaller than 0.03 is noise.\n\n"
    "v27: tried using a 2-patient validation holdout instead of LOPO. The validation "
    "AUC was so noisy on N=12 patients that the model picked epoch 0 as best, which "
    "is essentially the untrained initialization.\n\n"
    "v30: the L4 lecture's discriminative-LR recipe — backbone at 1/10 head LR — "
    "hurt us by a full 0.10 AUC. The reason: ImageNet features don't transfer "
    "cleanly to microscopy textures, so the backbone needs MORE learning rate, "
    "not less."
))


# =========================
# Slide 8: BACKBONE TESTS
# =========================
s = pres.slides.add_slide(BLANK)
set_bg(s, CREAM)
add_accent_bar(s)
add_title(s, "Clean backbone tests",
          eyebrow="Act 4  ·  the controls")

# Three columns: v34, v35, v36
def backbone_card(slide, x, y, w, h, name, backbone, lb, lb_color, blurb):
    add_rect(slide, x, y, w, h, WHITE)
    add_rect(slide, x, y, Inches(0.08), h, lb_color)
    add_text(slide, x + Inches(0.25), y + Inches(0.20), w - Inches(0.4), Inches(0.4),
             name, font="Georgia", size=22, bold=True, color=DARK,
             align=PP_ALIGN.LEFT)
    add_text(slide, x + Inches(0.25), y + Inches(0.70), w - Inches(0.4), Inches(0.35),
             backbone, font="Calibri", size=11, color=MUTED, italic=True,
             align=PP_ALIGN.LEFT)
    add_text(slide, x + Inches(0.25), y + Inches(1.15), w - Inches(0.4), Inches(0.7),
             lb, font="Georgia", size=32, bold=True, color=lb_color,
             align=PP_ALIGN.LEFT)
    add_text(slide, x + Inches(0.25), y + Inches(2.05), w - Inches(0.4), Inches(0.9),
             blurb, font="Calibri", size=10, color=DARK,
             align=PP_ALIGN.LEFT)

cw, ch = 2.95, 3.15
gx = 0.10
xs = [0.4, 0.4 + cw + gx, 0.4 + 2 * (cw + gx)]
backbone_card(s, Inches(xs[0]), Inches(1.40), Inches(cw), Inches(ch),
              "v34", "ResNet-50, 128 native", "0.7155", DARK,
              "Teacher's recommended backbone, tested cleanly: no upscale, no "
              "disc-LR. Lands at v23's seed-spread band of v19. Solid, not a win.")
backbone_card(s, Inches(xs[1]), Inches(1.40), Inches(cw), Inches(ch),
              "v35", "DenseNet-201, 128 native", "queued", MUTED,
              "Deeper dense connectivity. Submitted, waiting on LB.")
backbone_card(s, Inches(xs[2]), Inches(1.40), Inches(cw), Inches(ch),
              "v36", "SE-ResNet-50 + 20 ep", "queued", MUTED,
              "Isolates whether channel attention closes the EffNet ↔ ResNet "
              "gap. EfficientNet has SE built in; plain ResNet does not.")

# Bottom strip
add_rect(s, Inches(0.4), Inches(4.70), Inches(9.2), Inches(0.45), NAVY)
add_text(s, Inches(0.6), Inches(4.75), Inches(9.0), Inches(0.35),
         "Resolved by v34: v21's 0.7018 was the 224 upscale, not ResNet-50.",
         font="Georgia", size=12, italic=True, color=CREAM,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)

add_footer(s, 8)
set_notes(s, (
    "Act 4 final piece: clean backbone tests. We wanted to know whether EfficientNet-B0 "
    "was actually special, or whether v19 just got lucky.\n\n"
    "v34 is the answer. Same recipe as v19 — same augmentation, same MIL, same TTA "
    "— but the backbone swapped to ResNet-50. No upscale, no discriminative LR. "
    "Crucially, this is a single change from v19. LB came in at 0.7155.\n\n"
    "Two things to take from that. First, ResNet-50 is fine — it's not worse than "
    "EfficientNet in any meaningful way. The 0.7155 number sits inside the seed-"
    "variance band of v19 (0.7154 with seed 2 versus 0.7455 with seed 1). So we can "
    "now say: v21's earlier 0.7018 result was the 224 upscale's fault, not the "
    "backbone's. The 'no interpolation' warning from Lu 2020 was right.\n\n"
    "v35 with DenseNet-201 and v36 with SE-ResNet-50 are queued — both are tests "
    "of whether more parameters or channel attention closes the small remaining "
    "gap to v19's 0.7455. Results pending."
))


# =========================
# Slide 9: LEARNING CURVES
# =========================
s = pres.slides.add_slide(BLANK)
set_bg(s, CREAM)
add_accent_bar(s)
add_title(s, "Learning curves — v19 (LB 0.7455)",
          eyebrow="Evidence")

# Embed the v19 learning curves image
if LEARNING_CURVES_PNG.exists():
    # Available area: x=0.5..6.5 (w=6), y=1.4..4.8 (h=3.4)
    s.shapes.add_picture(str(LEARNING_CURVES_PNG),
                         Inches(0.5), Inches(1.4),
                         width=Inches(6.0))
else:
    add_rect(s, Inches(0.5), Inches(1.4), Inches(6.0), Inches(3.4), WHITE)
    add_text(s, Inches(0.5), Inches(2.8), Inches(6.0), Inches(0.5),
             "[learning curve image not found]",
             font="Calibri", size=12, italic=True, color=MUTED,
             align=PP_ALIGN.CENTER)

# Right: annotations
add_text(s, Inches(6.9), Inches(1.40), Inches(3.0), Inches(0.35),
         "What to notice", font="Calibri", size=12, bold=True, color=TERRA)
add_bullets(s, Inches(6.9), Inches(1.78), Inches(3.0), Inches(3.5), [
    ("Train and val track each other through training",),
    ("No early peak → no OOD overfit signal at the patient level seen during CV",),
    ("Stable convergence at ep ≈10",),
    ("Yet the LB sits 0.04 below CV — the remaining OOD gap that LOPO cannot expose",),
], size=11)

add_footer(s, 9)
set_notes(s, (
    "Learning curves from v19, the LB winner. A few things worth pointing out.\n\n"
    "First, the train and validation curves track each other through the run — no "
    "obvious overfitting signature, no early peak collapse like we saw in v15.\n\n"
    "Second, convergence is clean. The model stabilizes around epoch 10 and stays there.\n\n"
    "Third — and this is the assignment's take-home in one chart — the validation "
    "AUC looks great, but the public LB still sits about 0.04 below it. That's the "
    "OOD gap we cannot close from training-set evidence alone. The curves you can "
    "draw from training data have a floor below which they cannot warn you."
))


# =========================
# Slide 10: FINAL STRATEGY + LB POSITION
# =========================
s = pres.slides.add_slide(BLANK)
set_bg(s, CREAM)
add_accent_bar(s)
add_title(s, "Final-submission strategy", eyebrow="The plan")

# Left: 2-pick strategy
add_text(s, Inches(0.4), Inches(1.40), Inches(5.0), Inches(0.35),
         "Kaggle lets us pick 2 for private LB",
         font="Calibri", size=12, bold=True, color=TERRA)

# Pick 1 card
add_rect(s, Inches(0.4), Inches(1.80), Inches(5.4), Inches(1.30), WHITE)
add_rect(s, Inches(0.4), Inches(1.80), Inches(0.08), Inches(1.30), NAVY)
add_text(s, Inches(0.6), Inches(1.92), Inches(0.7), Inches(0.4),
         "01", font="Georgia", size=22, bold=True, color=TERRA)
add_text(s, Inches(1.35), Inches(1.92), Inches(4.0), Inches(0.4),
         "v19 — the safety net",
         font="Calibri", size=14, bold=True, color=DARK)
add_text(s, Inches(0.6), Inches(2.40), Inches(5.1), Inches(0.7),
         "Known public LB 0.7455. Our best confirmed single model. "
         "We don't gamble pick 1.",
         font="Calibri", size=11, color=MUTED, align=PP_ALIGN.LEFT)

# Pick 2 card
add_rect(s, Inches(0.4), Inches(3.25), Inches(5.4), Inches(1.30), WHITE)
add_rect(s, Inches(0.4), Inches(3.25), Inches(0.08), Inches(1.30), TERRA)
add_text(s, Inches(0.6), Inches(3.37), Inches(0.7), Inches(0.4),
         "02", font="Georgia", size=22, bold=True, color=TERRA)
add_text(s, Inches(1.35), Inches(3.37), Inches(4.0), Inches(0.4),
         "Best of {v34, v35, v36}",
         font="Calibri", size=14, bold=True, color=DARK)
add_text(s, Inches(0.6), Inches(3.85), Inches(5.1), Inches(0.7),
         "Clean backbone-comparison shot. If all three underperform v19, "
         "we pick v19 twice and still report the experiment honestly.",
         font="Calibri", size=11, color=MUTED, align=PP_ALIGN.LEFT)

# Right: leaderboard position visual
add_text(s, Inches(6.0), Inches(1.40), Inches(3.6), Inches(0.35),
         "LB position", font="Calibri", size=12, bold=True, color=TERRA)

# Vertical AUC scale: 0.50 at bottom, 0.80 at top
SCALE_X = Inches(6.4)
SCALE_TOP = Inches(1.85)
SCALE_H = Inches(2.7)
SCALE_W = Inches(0.4)
# Background track
add_rect(s, SCALE_X, SCALE_TOP, SCALE_W, SCALE_H, WHITE)
# Helper: y for an AUC (0.50 -> bottom, 0.80 -> top)
def auc_y(auc):
    frac = (0.80 - auc) / (0.80 - 0.50)
    return SCALE_TOP + Emu(int(SCALE_H * frac))

# Mark leader (0.7832)
y = auc_y(0.7832)
add_rect(s, SCALE_X - Inches(0.06), y - Inches(0.025), Inches(0.52), Inches(0.05), NAVY)
add_text(s, SCALE_X + Inches(0.55), y - Inches(0.18), Inches(3.0), Inches(0.32),
         "0.7832  leader",
         font="Calibri", size=11, bold=True, color=NAVY, align=PP_ALIGN.LEFT)

# Mark v19 (0.7455)
y = auc_y(0.7455)
add_rect(s, SCALE_X - Inches(0.06), y - Inches(0.025), Inches(0.52), Inches(0.05), TERRA)
add_text(s, SCALE_X + Inches(0.55), y - Inches(0.18), Inches(3.0), Inches(0.32),
         "0.7455  v19 (ours)",
         font="Calibri", size=11, bold=True, color=TERRA, align=PP_ALIGN.LEFT)

# Mark v34 (0.7155)
y = auc_y(0.7155)
add_rect(s, SCALE_X - Inches(0.06), y - Inches(0.025), Inches(0.52), Inches(0.05), SLATE)
add_text(s, SCALE_X + Inches(0.55), y - Inches(0.18), Inches(3.0), Inches(0.32),
         "0.7155  v34",
         font="Calibri", size=11, bold=True, color=SLATE, align=PP_ALIGN.LEFT)

# Mark v16 (0.503)  — wait, 0.503 is below scale (0.50 floor). Place at floor.
# Show as below-scale anchor
add_text(s, SCALE_X + Inches(0.55), SCALE_TOP + SCALE_H + Inches(0.05),
         Inches(3.0), Inches(0.32),
         "0.503  v16 (random baseline)",
         font="Calibri", size=10, italic=True, color=MUTED, align=PP_ALIGN.LEFT)

# Scale labels
add_text(s, SCALE_X - Inches(0.55), SCALE_TOP - Inches(0.08), Inches(0.5), Inches(0.3),
         "0.80", font="Calibri", size=8, color=MUTED, align=PP_ALIGN.RIGHT)
add_text(s, SCALE_X - Inches(0.55), SCALE_TOP + SCALE_H - Inches(0.10),
         Inches(0.5), Inches(0.3),
         "0.50", font="Calibri", size=8, color=MUTED, align=PP_ALIGN.RIGHT)

add_footer(s, 10)
set_notes(s, (
    "Kaggle lets you pick two submissions for the private leaderboard. Our plan:\n\n"
    "Pick 1 is v19 — our known safety net at 0.7455. We don't gamble this slot.\n\n"
    "Pick 2 is the best of the backbone-comparison runs — whichever of v34, v35, v36 "
    "lands highest by competition close on June 3. If all three underperform v19, "
    "we pick v19 twice and our report still tells the honest experimental story.\n\n"
    "Where we sit relative to the leaderboard: leader is 0.7832, we're at 0.7455, "
    "and v34 is at 0.7155. The gap to leader is 0.038. v16 — our heaviest engineering "
    "effort — is shown below the scale at 0.503, basically random. The chart is a "
    "compact picture of the whole story."
))


# =========================
# Slide 11: ACT 5 — THE SEMI-SUPERVISED PIVOT (post-deck-build addition)
# =========================
s = pres.slides.add_slide(BLANK)
set_bg(s, CREAM)
add_accent_bar(s)
add_title(s, "Act 5: the semi-supervised pivot took us to #1",
          eyebrow="Postscript  ·  v41 → v47")

# LB chart on the right side
LB_CHART_PNG = PROJECT_ROOT / "report" / "figures" / "lb_progression.png"
if LB_CHART_PNG.exists():
    s.shapes.add_picture(str(LB_CHART_PNG), Inches(4.2), Inches(1.5),
                         width=Inches(5.6), height=Inches(2.9))

# Left-column bullets summarizing the arc
add_bullets(s, Inches(0.4), Inches(1.5), Inches(3.7), Inches(3.6), [
    ("v41: +0.011 LB from 4 L4 regularizers",),
    ("v43: stacked 4 more changes → −0.012 regression",),
    ("v44: hard pseudo (Lee 2013) +0.025 LB",),
    ("v46: SOFT pseudo / Hinton 2015 distillation +0.039 LB",),
    ("v47: noisy-student round 2 (Xie 2020) +0.003 LB",),
    ("Result: 0.8264, #1 on public LB (+0.013 lead)",),
], size=12)

# Bottom banner — the headline
add_rect(s, Inches(0.4), Inches(4.55), Inches(9.2), Inches(0.55), NAVY)
add_text(s, Inches(0.6), Inches(4.62), Inches(9.0), Inches(0.42),
         "Dataset-size lever (+0.067 from pseudo+distillation) dwarfed regularization gains (+0.011 from v41).",
         font="Georgia", size=13, italic=True, color=CREAM,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)

add_footer(s, 11)
set_notes(s, (
    "Act 5 — the chapter that didn't exist when this deck was first built.\n\n"
    "Quick arc: v41 stacked four L4-textbook regularizers cleanly for +0.011 LB. v43 then "
    "stacked four MORE changes simultaneously and regressed by 0.012 — a methodological "
    "warning we paid the price to learn.\n\n"
    "v44 was the first breakthrough: Lee 2013 pseudo-labels. Teacher (v41) predicts on the "
    "test set; cells with confidence >0.95 or <0.05 get added to training with patient_id=-1 "
    "so MIL skips them. +0.025 LB.\n\n"
    "v46 was the bigger breakthrough: Hinton 2015 distillation. Use ALL 59,000 test cells, "
    "not just the confident ones. Soft target = teacher's raw probability. +0.039 LB. We "
    "took #1 on the public LB at 0.8236.\n\n"
    "v47 closed the loop: Xie 2020-style iterative noisy student round 2. Same recipe as v46 "
    "but with v46's ensemble as the new teacher. +0.003 LB — small, but enough to extend "
    "our LB lead to +0.013 over next-best. The per-seed tr_auc spread also tightened (0.989-0.993 "
    "in v46 → 0.993-0.994 in v47), so round 2 acts as a variance reducer as well.\n\n"
    "The key methodological insight (bottom banner): the +0.067 LB lift from pseudo+distillation "
    "dwarfed the +0.011 from textbook regularization. On 12 patients, the dataset-size lever "
    "is bigger than the model-quality lever."
))


# =========================
# Slide 12: TAKE-HOME / CLOSING
# =========================
s = pres.slides.add_slide(BLANK)
set_bg(s, NAVY)

# Terracotta bar on left
add_rect(s, Inches(0), Inches(0), Inches(0.25), Inches(5.625), TERRA)

# Eyebrow + title
add_text(s, Inches(0.8), Inches(0.45), Inches(8), Inches(0.3),
         "FIVE THINGS WE'LL CARRY FORWARD",
         font="Calibri", size=12, bold=True, color=TERRA, align=PP_ALIGN.LEFT)
add_text(s, Inches(0.8), Inches(0.85), Inches(8.5), Inches(0.7),
         "Take-home messages", font="Georgia", size=32, bold=True,
         color=WHITE, align=PP_ALIGN.LEFT)

takeaways = [
    ("Higher CV ≠ higher LB.",
     "v15 had the best CV we ever saw (0.866). It had the worst LB (0.572)."),
    ("Don't stack unvalidated changes.",
     "v43 added 4 simultaneous tweaks on top of v41 → −0.012 LB regression with no way to attribute blame."),
    ("Seed variance is huge.",
     "Same recipe, different seed: 0.03 LB shift. Most single-seed A/B claims < 0.03 are noise."),
    ("On small-N, dataset size > regularization.",
     "v41 stacked four regularizers: +0.011 LB. v46 added soft pseudo-labels for ~50k extra cells: +0.039 LB."),
    ("Negative results are the methodology.",
     "v42 SSL collapse, v22 + v45_probe ensemble failures, v43 regression — every failure shaped the next version."),
]
y0 = 1.85
for i, (head, body) in enumerate(takeaways):
    y = Inches(y0 + i * 0.62)
    # Number
    add_text(s, Inches(0.8), y, Inches(0.4), Inches(0.35),
             f"{i+1:02d}", font="Georgia", size=15, bold=True, color=TERRA,
             align=PP_ALIGN.LEFT)
    # Headline + body inline
    add_text(s, Inches(1.3), y, Inches(8.2), Inches(0.32),
             head, font="Calibri", size=13, bold=True, color=WHITE,
             align=PP_ALIGN.LEFT)
    add_text(s, Inches(1.3), y + Inches(0.30), Inches(8.2), Inches(0.32),
             body, font="Calibri", size=11, italic=True, color=CREAM,
             align=PP_ALIGN.LEFT)

# Footer
add_text(s, Inches(0.8), Inches(5.30), Inches(8), Inches(0.25),
         "Thanks. Code + LB history at github.com/rafallex/A3-ADL  ·  Questions?",
         font="Calibri", size=10, italic=True, color=MUTED, align=PP_ALIGN.LEFT)

set_notes(s, (
    "Five things we'll carry forward. I'll read them quickly and leave the slide up for "
    "questions.\n\n"
    "1. Higher CV doesn't mean higher LB. v15 had our best CV — 0.866 — and our worst "
    "LB — 0.572. If you only have CV, you don't know what you have.\n\n"
    "2. Don't stack unvalidated changes. v43 added four tweaks at once on top of v41 and "
    "regressed by 0.012 LB. We had no way to attribute blame and the run cost 5 hours of "
    "Kaggle GPU. After v43 we adopted a one-or-two-variables-per-experiment rule.\n\n"
    "3. Seed variance is huge. v23 was v19 with a different random seed. 0.03 LB shift. "
    "Most A/B comparisons in the literature with deltas under 0.03 are not reliable.\n\n"
    "4. On small-N patient-grouped data, dataset size beats regularization. v41 stacked "
    "four textbook regularizers for +0.011 LB. v46 added soft pseudo-labels for ~50,000 "
    "extra effective cells: +0.039 LB. Three and a half times the gain from a different "
    "lever entirely.\n\n"
    "5. Negative results are the methodology. v42 SSL collapse, v22 and v45_probe ensemble "
    "failures, v43 stacked regression — six diagnosed failure modes, each shaped the next "
    "version. The wins came from cleanly understanding the failures.\n\n"
    "Thanks. Code + full LB history at github.com/rafallex/A3-ADL. Questions?"
))


# --- Save ---
pres.save(str(OUT_PATH))
print(f"Saved: {OUT_PATH}")
