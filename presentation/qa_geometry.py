"""Analytical QA for A3_cancer_challenge.pptx.

Checks:
1. Every shape stays inside slide bounds (10" x 5.625")
2. For each slide, detect pairs of TEXTBOXES whose bounding boxes overlap
   (overlap of textboxes with rectangles is usually intentional — text on a card;
   overlap of textboxes with textboxes is usually a layout bug).

Reports as a punch-list of (slide, shape, issue).
"""
from pptx import Presentation
from pptx.util import Emu

SLIDE_W_EMU = 9144000   # 10"
SLIDE_H_EMU = 5143500   # 5.625"
TOLERANCE_EMU = Emu(0.02 * 914400)  # 0.02" slack


def emu_to_in(v):
    return v / 914400.0


def main():
    p = Presentation("A3_cancer_challenge.pptx")
    issues = []

    for i, slide in enumerate(p.slides, 1):
        # Collect textbox rectangles (and labels)
        textboxes = []
        for shape in slide.shapes:
            x, y, w, h = shape.left, shape.top, shape.width, shape.height
            if x is None:
                continue

            # 1. Slide-bounds check
            if x < -TOLERANCE_EMU:
                issues.append((i, shape.shape_id, f"left x={emu_to_in(x):.2f} < 0"))
            if y < -TOLERANCE_EMU:
                issues.append((i, shape.shape_id, f"top y={emu_to_in(y):.2f} < 0"))
            if x + w > SLIDE_W_EMU + TOLERANCE_EMU:
                issues.append((i, shape.shape_id,
                              f"overflows right: x+w={emu_to_in(x+w):.2f} > 10.00"))
            if y + h > SLIDE_H_EMU + TOLERANCE_EMU:
                issues.append((i, shape.shape_id,
                              f"overflows bottom: y+h={emu_to_in(y+h):.3f} > 5.625"))

            # Collect textbox-only rectangles for overlap check
            if shape.has_text_frame and shape.shape_type is not None:
                # python-pptx: shape_type 17 = MSO_SHAPE_TYPE.TEXT_BOX
                if int(shape.shape_type) == 17:
                    txt_snip = "".join(
                        r.text for p_ in shape.text_frame.paragraphs for r in p_.runs
                    )[:40]
                    textboxes.append((shape.shape_id, x, y, w, h, txt_snip))

        # 2. Textbox-textbox overlap
        for j in range(len(textboxes)):
            for k in range(j + 1, len(textboxes)):
                sid1, x1, y1, w1, h1, t1 = textboxes[j]
                sid2, x2, y2, w2, h2, t2 = textboxes[k]
                # Bounding-box overlap (skip near-zero overlaps)
                overlap_x = min(x1 + w1, x2 + w2) - max(x1, x2)
                overlap_y = min(y1 + h1, y2 + h2) - max(y1, y2)
                if overlap_x > TOLERANCE_EMU and overlap_y > TOLERANCE_EMU:
                    issues.append((i, f"{sid1}↔{sid2}",
                                   f"textbox overlap "
                                   f"({emu_to_in(overlap_x):.2f}×{emu_to_in(overlap_y):.2f}\"): "
                                   f"\"{t1}\" vs \"{t2}\""))

    if not issues:
        print("No geometry issues found.")
        return

    print(f"Found {len(issues)} potential issues:\n")
    # Group by slide
    from collections import defaultdict
    by_slide = defaultdict(list)
    for slide_idx, sid, msg in issues:
        by_slide[slide_idx].append((sid, msg))
    for slide_idx in sorted(by_slide):
        print(f"=== Slide {slide_idx} ===")
        for sid, msg in by_slide[slide_idx]:
            print(f"  [{sid}] {msg}")
        print()


if __name__ == "__main__":
    main()
